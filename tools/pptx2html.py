#!/usr/bin/env python3
# Conversor PPTX -> HTML de alta fidelidad (formas absolutas a 1280x720).
import zipfile, os, html, hashlib, sys, json, io
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.dml import MSO_THEME_COLOR, MSO_COLOR_TYPE
from lxml import etree

NS = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
      'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}

THEME_NAME = {
    MSO_THEME_COLOR.ACCENT_1: 'accent1', MSO_THEME_COLOR.ACCENT_2: 'accent2',
    MSO_THEME_COLOR.ACCENT_3: 'accent3', MSO_THEME_COLOR.ACCENT_4: 'accent4',
    MSO_THEME_COLOR.ACCENT_5: 'accent5', MSO_THEME_COLOR.ACCENT_6: 'accent6',
    MSO_THEME_COLOR.DARK_1: 'dk1', MSO_THEME_COLOR.LIGHT_1: 'lt1',
    MSO_THEME_COLOR.DARK_2: 'dk2', MSO_THEME_COLOR.LIGHT_2: 'lt2',
    MSO_THEME_COLOR.TEXT_1: 'tx1', MSO_THEME_COLOR.TEXT_2: 'tx2',
    MSO_THEME_COLOR.BACKGROUND_1: 'bg1', MSO_THEME_COLOR.BACKGROUND_2: 'bg2',
    MSO_THEME_COLOR.HYPERLINK: 'hlink', MSO_THEME_COLOR.FOLLOWED_HYPERLINK: 'folHlink',
}


def parse_theme(zf):
    scheme = {}; clrmap = {}
    try:
        root = etree.fromstring(zf.read('ppt/theme/theme1.xml'))
        cs = root.find('.//a:clrScheme', NS)
        for child in cs:
            tag = etree.QName(child).localname
            srgb = child.find('a:srgbClr', NS); sysc = child.find('a:sysClr', NS)
            if srgb is not None:
                scheme[tag] = '#' + srgb.get('val')
            elif sysc is not None:
                scheme[tag] = '#' + (sysc.get('lastClr') or '000000')
    except Exception:
        pass
    try:
        for name in zf.namelist():
            if name.startswith('ppt/slideMasters/slideMaster') and name.endswith('.xml'):
                mroot = etree.fromstring(zf.read(name))
                cm = mroot.find('.//p:clrMap', NS)
                if cm is not None:
                    clrmap = dict(cm.attrib)
                break
    except Exception:
        pass
    return scheme, clrmap

def repair(src, dst):
    raw = zipfile.ZipFile(src); zo = zipfile.ZipFile(dst, 'w', zipfile.ZIP_DEFLATED)
    for it in raw.infolist():
        if it.filename.lower().endswith('.fntdata'):
            continue
        try:
            zo.writestr(it, raw.read(it.filename))
        except Exception:
            pass
    zo.close()

def esc(t):
    return html.escape(t or '', quote=True)

ALIGN = {PP_ALIGN.CENTER: 'center', PP_ALIGN.RIGHT: 'right', PP_ALIGN.JUSTIFY: 'justify', PP_ALIGN.LEFT: 'left'}
ANCHOR = {MSO_ANCHOR.MIDDLE: 'center', MSO_ANCHOR.BOTTOM: 'flex-end', MSO_ANCHOR.TOP: 'flex-start'}

class Conv:
    def __init__(self, prs, imgdir, prefix, pxperemu_x, pxperemu_y, pxperpt, scheme=None, clrmap=None):
        self.prs = prs; self.imgdir = imgdir; self.prefix = prefix
        self.kx = pxperemu_x; self.ky = pxperemu_y; self.pxperpt = pxperpt
        self.imgcount = 0
        self.bg = []  # elementos de fondo en z-order: (bounds_px, 'solid'|'pic', color_hex|PIL) para muestrear color bajo el texto
        self.scheme = scheme or {}; self.clrmap = clrmap or {}
        self.steps = {}  # id(elemento) -> índice de paso de animación (solo TEXTO, ordenado arriba->abajo)
        self.gfx = {}    # id(elemento) -> índice de gráfico (fotos/ilustraciones/paneles: entran al inicio)
        self.imgweb = 'img/pptx/'  # carpeta destino de imágenes en el sitio

    def is_content(self, sp, b):
        x, y, w, h = b
        if (w * h) / (1280 * 720) > 0.55:
            return False
        if w < 40 or h < 20:
            return False
        if sp.shape_type == MSO_SHAPE_TYPE.PICTURE:
            return True
        try:
            if sp.has_text_frame and sp.text_frame.text.strip():
                return True
        except Exception:
            pass
        try:
            if self.solid_color(sp.fill):
                return True
        except Exception:
            pass
        return False

    def plan_steps(self, shapes):
        # Separa el contenido en dos grupos:
        #  - gfx: fotos / ilustraciones (grupos) / paneles de color sin texto -> entran al inicio de la diapositiva
        #  - steps: bloques de TEXTO -> se revelan siguiendo la narración
        txt, gfx = [], []
        for sp in shapes:
            try:
                b = self.px_bounds(sp, (1.0, 0.0, 1.0, 0.0))
                if not b:
                    continue
                eid = sp.shape_id
                if sp.shape_type == MSO_SHAPE_TYPE.GROUP:
                    x, y, w, h = b
                    if (w * h) / (1280 * 720) < 0.55 and w > 40 and h > 25:
                        gfx.append((y, x, eid))
                    continue
                if not self.is_content(sp, b):
                    continue
                has_txt = False
                try:
                    has_txt = sp.has_text_frame and sp.text_frame.text.strip() != ''
                except Exception:
                    pass
                data, _ = self.shape_image(sp)
                if data or not has_txt:
                    gfx.append((b[1], b[0], eid))   # imagen (emit la trata como <img>) o panel sin texto
                else:
                    txt.append((b[1], b[0], eid))
            except Exception:
                pass
        txt.sort(key=lambda t: (round(t[0] / 42), t[1]))
        gfx.sort(key=lambda t: (round(t[0] / 42), t[1]))
        self.steps = {eid: i for i, (y, x, eid) in enumerate(txt)}
        self.gfx = {eid: i for i, (y, x, eid) in enumerate(gfx)}

    def resolve(self, name):
        if name in self.clrmap:
            name = self.clrmap[name]
        return self.scheme.get(name)

    def _hex(self, cf):
        try:
            if cf is None:
                return None
            if cf.type == MSO_COLOR_TYPE.RGB:
                rgb = cf.rgb
                return '#%02x%02x%02x' % (rgb[0], rgb[1], rgb[2])
            if cf.type == MSO_COLOR_TYPE.SCHEME:
                nm = THEME_NAME.get(cf.theme_color)
                if nm:
                    return self.resolve(nm)
        except Exception:
            pass
        return None

    def px_bounds(self, sp, tf):
        ax, bx, ay, by = tf
        try:
            L = sp.left; T = sp.top; W = sp.width; H = sp.height
        except Exception:
            return None
        if L is None:
            return None
        return ((ax * L + bx) * self.kx, (ay * T + by) * self.ky, ax * W * self.kx, ay * H * self.ky)

    @staticmethod
    def is_dark(hexc):
        r = int(hexc[1:3], 16); g = int(hexc[3:5], 16); b = int(hexc[5:7], 16)
        return (0.299 * r + 0.587 * g + 0.114 * b) < 130

    def collect_bg(self, shapes, tf):
        for sp in shapes:
            try:
                if sp.shape_type == MSO_SHAPE_TYPE.GROUP:
                    ntf = self.group_tf(sp, tf)
                    if ntf:
                        self.collect_bg(sp.shapes, ntf)
                    continue
                b = self.px_bounds(sp, tf)
                if not b:
                    continue
                if sp.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        im = Image.open(io.BytesIO(sp.image.blob)).convert('RGB')
                        im.thumbnail((160, 160))
                        self.bg.append((b, 'pic', im))
                    except Exception:
                        pass
                    continue
                col = self.solid_color(sp.fill) if hasattr(sp, 'fill') else None
                if col:
                    self.bg.append((b, 'solid', col))
            except Exception:
                pass

    def group_tf(self, sp, tf):
        ax, bx, ay, by = tf
        try:
            off = sp._element.find('.//a:off', NS); ext = sp._element.find('.//a:ext', NS)
            choff = sp._element.find('.//a:chOff', NS); chext = sp._element.find('.//a:chExt', NS)
            gL = int(off.get('x')); gT = int(off.get('y')); gW = int(ext.get('cx')); gH = int(ext.get('cy'))
            coX = int(choff.get('x')); coY = int(choff.get('y')); ceW = int(chext.get('cx')); ceH = int(chext.get('cy'))
            nax = ax * (gW / ceW); nbx = ax * (gL - coX * (gW / ceW)) + bx
            nay = ay * (gH / ceH); nby = ay * (gT - coY * (gH / ceH)) + by
            return (nax, nbx, nay, nby)
        except Exception:
            return None

    def over_dark(self, bounds):
        if not bounds:
            return False
        cx = bounds[0] + bounds[2] / 2; cy = bounds[1] + bounds[3] / 2
        # el elemento de fondo más "arriba" (último en z-order) que contiene el centro
        for (b, kind, data) in reversed(self.bg):
            x, y, w, h = b
            if not (x <= cx <= x + w and y <= cy <= y + h):
                continue
            if w <= 0 or h <= 0:
                continue
            if kind == 'solid':
                return self.is_dark(data)
            try:
                iw, ih = data.size
                sx = int(min(iw - 1, max(0, (cx - x) / w * iw)))
                sy = int(min(ih - 1, max(0, (cy - y) / h * ih)))
                r, g, bl = data.getpixel((sx, sy))[:3]
                return (0.299 * r + 0.587 * g + 0.114 * bl) < 130
            except Exception:
                return False
        return False

    def shape_image(self, sp):
        # extrae la imagen de una forma: Picture, placeholder de imagen, o relleno de imagen (blipFill)
        try:
            img = sp.image
            return img.blob, img.ext
        except Exception:
            pass
        try:
            blip = sp._element.find('.//a:blip', NS)
            if blip is not None:
                rId = blip.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
                if rId:
                    part = sp.part.related_part(rId)
                    ext = (part.content_type or 'image/png').split('/')[-1]
                    return part.blob, ext
        except Exception:
            pass
        return None, None

    def solid_color(self, fill):
        try:
            if fill.type == 1:  # SOLID
                return self._hex(fill.fore_color)
        except Exception:
            pass
        return None

    def run_color(self, run):
        try:
            return self._hex(run.font.color)
        except Exception:
            pass
        return None

    def para_html(self, p, default_color):
        runs = []
        for r in p.runs:
            t = esc(r.text)
            if not t:
                continue
            st = []
            sz = None
            try:
                sz = r.font.size.pt if r.font.size else None
            except Exception:
                sz = None
            if sz:
                st.append('font-size:%.1fpx' % (sz * self.pxperpt))
            if r.font.bold:
                st.append('font-weight:700')
            if r.font.italic:
                st.append('font-style:italic')
            col = self.run_color(r) or default_color
            if col:
                st.append('color:' + col)
            runs.append('<span style="%s">%s</span>' % (';'.join(st), t))
        if not runs:
            return ''
        align = ALIGN.get(p.alignment, 'left')
        lvl = (p.level or 0)
        indent = ('margin-left:%dpx;' % (lvl * 18)) if lvl else ''
        return '<div style="text-align:%s;%s">%s</div>' % (align, indent, ''.join(runs))

    def text_html(self, shape, force_white=False):
        tf = shape.text_frame
        default_color = '#fff' if force_white else '#000'
        paras = [self.para_html(p, default_color) for p in tf.paragraphs]
        paras = [x for x in paras if x]
        if not paras:
            return None
        anchor = 'flex-start'
        try:
            anchor = ANCHOR.get(tf.vertical_anchor, 'flex-start')
        except Exception:
            pass
        inner = ''.join(paras)
        return ('<div style="display:flex;flex-direction:column;justify-content:%s;'
                'height:100%%;width:100%%;line-height:1.15;overflow:hidden">%s</div>' % (anchor, inner))

    def emit_shape(self, sp, tf, out, steppable=True):
        # tf: (ax,bx,ay,by) mapea coords de este nivel -> EMU de slide
        ax, bx, ay, by = tf
        try:
            L = sp.left; T = sp.top; W = sp.width; H = sp.height
        except Exception:
            L = T = W = H = None
        if L is None:
            return
        # a EMU de slide
        sL = ax * L + bx; sT = ay * T + by
        sW = ax * W; sH = ay * H
        # a px
        px = sL * self.kx; py = sT * self.ky
        pw = sW * self.kx; ph = sH * self.ky
        rot = 0
        try:
            rot = sp.rotation or 0
        except Exception:
            rot = 0
        base = 'position:absolute;left:%.1fpx;top:%.1fpx;width:%.1fpx;height:%.1fpx;' % (px, py, pw, ph)
        if rot:
            base += 'transform:rotate(%.2fdeg);' % rot
        try:
            eid = sp.shape_id if steppable else None
        except Exception:
            eid = None
        ds = (' data-step="%d"' % self.steps[eid]) if eid in self.steps else ''
        gi = self.gfx.get(eid) if eid is not None else None
        gcls = ' class="gfx"' if gi is not None else ''
        gdel = ('animation-delay:%.2fs;' % min(gi * 0.13, 0.78)) if gi is not None else ''

        st = sp.shape_type
        if st == MSO_SHAPE_TYPE.GROUP:
            ntf = self.group_tf(sp, tf)
            if ntf:
                if ds or gcls:
                    out.append('<div%s%s style="position:absolute;inset:0;pointer-events:none;%s">' % (ds, gcls, gdel))
                for child in sp.shapes:
                    self.emit_shape(child, ntf, out, steppable=False)
                if ds or gcls:
                    out.append('</div>')
            return

        data, ext = self.shape_image(sp)
        if data:
            try:
                h = hashlib.md5(data).hexdigest()[:10]
                fn = '%s_%s.%s' % (self.prefix, h, ext)
                fp = os.path.join(self.imgdir, fn)
                if not os.path.exists(fp):
                    open(fp, 'wb').write(data)
                area = (pw * ph) / (1280 * 720)
                fit = 'cover' if area > 0.15 else 'contain'
                out.append('<img%s%s src="%s%s" style="%s%soverflow:hidden;object-fit:%s" alt="">' % (ds, gcls, self.imgweb, fn, base, gdel, fit))
            except Exception:
                pass
            return

        # autoshape / textbox
        color = self.solid_color(sp.fill) if hasattr(sp, 'fill') else None
        radius = ''
        try:
            if 'roundRect' in (sp._element.xml or ''):
                radius = 'border-radius:14px;'
        except Exception:
            pass
        has_text = False
        try:
            has_text = sp.has_text_frame and sp.text_frame.text.strip() != ''
        except Exception:
            has_text = False
        bg = ('background:%s;' % color) if color else ''
        pad = 'padding:6px 10px;' if has_text else ''
        if color or has_text:
            inner = ''
            if has_text:
                fw = (color is not None and self.is_dark(color)) or self.over_dark((px, py, pw, ph))
                inner = self.text_html(sp, fw) or ''
            out.append('<div%s%s style="%s%s%s%s%s">%s</div>' % (ds, gcls, base, bg, radius, pad, gdel, inner))

    def slide_html(self, slide):
        ident = (1.0, 0.0, 1.0, 0.0)
        self.bg = []
        layout = master = None
        try:
            layout = slide.slide_layout; master = layout.slide_master
            self.collect_bg([sp for sp in master.shapes if not sp.is_placeholder], ident)
            self.collect_bg([sp for sp in layout.shapes if not sp.is_placeholder], ident)
        except Exception:
            pass
        self.collect_bg(list(slide.shapes), ident)
        self.steps = {}
        self.plan_steps(list(slide.shapes))
        out = []
        try:
            if master is not None:
                for sp in master.shapes:
                    if not sp.is_placeholder:
                        self.emit_shape(sp, ident, out, steppable=False)
            if layout is not None:
                for sp in layout.shapes:
                    if not sp.is_placeholder:
                        self.emit_shape(sp, ident, out, steppable=False)
        except Exception:
            pass
        for sp in slide.shapes:
            self.emit_shape(sp, ident, out)
        return ''.join(out)


def convert(pptx_path, prefix, imgdir, outjson):
    scratch = os.path.dirname(outjson)
    rep = os.path.join(scratch, prefix + '_rep.pptx')
    repair(pptx_path, rep)
    prs = Presentation(rep)
    sw, sh = prs.slide_width, prs.slide_height
    kx = 1280.0 / sw; ky = 720.0 / sh
    pxperpt = 1280.0 * 12700.0 / sw
    os.makedirs(imgdir, exist_ok=True)
    scheme, clrmap = parse_theme(zipfile.ZipFile(rep))
    c = Conv(prs, imgdir, prefix, kx, ky, pxperpt, scheme, clrmap)
    slides = []
    for i, sl in enumerate(prs.slides, 1):
        slides.append(c.slide_html(sl))
    json.dump(slides, open(outjson, 'w'), ensure_ascii=False)
    print('%s: %d slides -> %s (imgs en %s)' % (prefix, len(slides), outjson, imgdir))
    return slides


if __name__ == '__main__':
    base = "/Users/jonathanf/Desktop/GRI Tutorias/2026 GRI TUTORIAS/Curso Blueprint"
    imgdir = base + "/curso-html5/img/pptx"
    scratch = "/private/tmp/claude-501/-Users-jonathanf-Desktop-App-Labsream/67bfb052-9191-4e04-8dbe-0e60b19b60d9/scratchpad"
    convert(base + "/Introducción/Introducción_vf.pptx", "intro", imgdir, scratch + "/intro_slides.json")
