#!/usr/bin/env python3
# Genera slides_html.json: <img> del PDF + "coberturas" que se revelan poco a poco
# (una por zona de contenido, detectadas desde el PPTX). Reutiliza pptx2html.Conv.
import sys, os, json, zipfile
sys.path.insert(0, os.path.dirname(__file__))
import pptx2html as P
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

BASE = "/Users/jonathanf/Desktop/GRI Tutorias/2026 GRI TUTORIAS/Curso Blueprint"
IMGDIR = BASE + "/curso-html5/img/slides"

def light(c):
    return c[0] > 232 and c[1] > 232 and c[2] > 232

def uniform_bg(im):
    w, h = im.size
    # rejilla de puntos evitando el borde; mayoría clara => fondo claro
    pts = []
    for gy in range(1, 8):
        for gx in range(1, 12):
            pts.append((int(gx / 12 * w), int(gy / 8 * h)))
    cs = [im.getpixel(p)[:3] for p in pts]
    lightones = [c for c in cs if light(c)]
    if len(lightones) < 0.62 * len(cs):   # no es fondo claro dominante
        return None
    r = sum(c[0] for c in lightones) // len(lightones)
    g = sum(c[1] for c in lightones) // len(lightones)
    b = sum(c[2] for c in lightones) // len(lightones)
    return '#%02x%02x%02x' % (r, g, b)

def regions_for(prs, conv, slide):
    regs = []
    def walk(shapes, tf):
        for sp in shapes:
            try:
                b = conv.px_bounds(sp, tf)
                st = sp.shape_type
                if st == MSO_SHAPE_TYPE.GROUP:
                    if b:
                        x, y, w, h = b
                        if (w * h) / (1280 * 720) < 0.55 and w > 45 and h > 30:
                            regs.append(b)          # el grupo (ilustración) como una zona
                    continue
                if not b:
                    continue
                x, y, w, h = b
                if (w * h) / (1280 * 720) > 0.55:
                    continue                        # fondo a sangre
                if w < 45 or h < 22:
                    continue                        # logo/nº de página/deco
                is_content = False
                try:
                    if sp.has_text_frame and sp.text_frame.text.strip():
                        is_content = True
                except Exception:
                    pass
                if st == MSO_SHAPE_TYPE.PICTURE:
                    is_content = True
                try:
                    if conv.solid_color(sp.fill):
                        is_content = True
                except Exception:
                    pass
                if is_content:
                    regs.append(b)
            except Exception:
                pass
    walk(list(slide.shapes), (1.0, 0.0, 1.0, 0.0))
    # ordenar arriba->abajo, izq->der
    regs.sort(key=lambda r: (round(r[1] / 45), r[0]))
    return regs

def convert(prefix, pptx_rel, id_prefix, n_slides, offset=0):
    rep = "/private/tmp/claude-501/-Users-jonathanf-Desktop-App-Labsream/67bfb052-9191-4e04-8dbe-0e60b19b60d9/scratchpad/%s_rep.pptx" % prefix
    P.repair(BASE + "/" + pptx_rel, rep)
    prs = Presentation(rep)
    sw, sh = prs.slide_width, prs.slide_height
    scheme, clrmap = P.parse_theme(zipfile.ZipFile(rep))
    conv = P.Conv(prs, "x", prefix, 1280.0 / sw, 720.0 / sh, 1280.0 * 12700 / sw, scheme, clrmap)
    slides = list(prs.slides)
    out = {}
    for i in range(1, n_slides + 1):
        sid = "%s%02d" % (id_prefix, i)
        img = "%s_%02d.png" % (prefix, i + offset)
        frag = '<img src="img/slides/%s" style="position:absolute;inset:0;width:100%%;height:100%%;object-fit:contain" alt="">' % img
        # detectar bg uniforme claro -> revelar por zonas
        ip = IMGDIR + "/" + img
        anim = ''
        try:
            im = Image.open(ip).convert('RGB')
            bg = uniform_bg(im)
            if bg:
                regs = regions_for(prs, conv, slides[i - 1 + offset])
                body = [r for r in regs if r[1] > 66]   # el título/zona superior queda estático
                for k, (x, y, w, h) in enumerate(body):
                    M = 9  # margen para que la tapa cubra por completo el contenido base (evita doblado)
                    L = max(0, x - M); T = max(0, y - M)
                    W = min(1280, x + w + M) - L; H = min(720, y + h + M) - T
                    if W <= 0 or H <= 0:
                        continue
                    # tapa permanente (oculta el contenido base) + recorte que entra con animación
                    anim += ('<div class="rev-hide" style="left:%.1fpx;top:%.1fpx;width:%.1fpx;height:%.1fpx;background:%s"></div>'
                             % (L, T, W, H, bg))
                    anim += ('<div class="rev-show" data-step="%d" style="left:%.1fpx;top:%.1fpx;width:%.1fpx;height:%.1fpx;'
                             'background-image:url(img/slides/%s);background-size:1280px 720px;background-position:-%.1fpx -%.1fpx;background-repeat:no-repeat"></div>'
                             % (k, L, T, W, H, img, L, T))
        except Exception:
            pass
        out[sid] = frag + anim
    return out

if __name__ == '__main__':
    allslides = {}
    allslides.update(convert("intro", "Introducción/Introducción_vf.pptx", "s", 9))
    allslides.update(convert("m1", "Módulo 1/Módulo 1_vf.pptx", "m1s", 14))
    json.dump(allslides, open(BASE + "/curso-html5/data/slides_html.json", "w"), ensure_ascii=False)
    ncov = sum(v.count('rev-cover') for v in allslides.values())
    print("slides:", len(allslides), "| coberturas de revelado:", ncov)
